#!/usr/bin/env python
# -*- coding: utf-8 -*-

import streamlit as st
import os
from pptx import Presentation  # Fixed import statement
from pathlib import Path
import sys

def convert_ppt_to_markdown(ppt_path):
    try:
        prs = Presentation(ppt_path)
        markdown_content = []
        
        for slide_number, slide in enumerate(prs.slides, 1):
            markdown_content.append(f"## Slide {slide_number}\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    markdown_content.append(shape.text.strip() + "\n")
            
            markdown_content.append("\n---\n")
        
        return "".join(markdown_content)
    except Exception as e:
        st.error(f"Error converting {ppt_path}: {str(e)}")
        return None

def save_markdown(content, output_path):
    try:
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(content)
    except Exception as e:
        st.error(f"Error saving markdown file {output_path}: {str(e)}")

def main():
    try:
        st.title("PowerPoint to Markdown Converter")
        
        folder_path = st.text_input("Enter folder path containing PPT files:")
        
        if folder_path and os.path.exists(folder_path):
            ppt_files = [f for f in os.listdir(folder_path) if f.endswith(('.ppt', '.pptx'))]
            
            if not ppt_files:
                st.warning("No PowerPoint files found in the specified folder.")
                return
            
            st.write(f"Found {len(ppt_files)} PowerPoint files")
            
            if st.button("Convert to Markdown"):
                output_folder = os.path.join(folder_path, 'markdown_output')
                os.makedirs(output_folder, exist_ok=True)
                
                for ppt_file in ppt_files:
                    ppt_path = os.path.join(folder_path, ppt_file)
                    markdown_content = convert_ppt_to_markdown(ppt_path)
                    
                    if markdown_content:
                        output_file = os.path.join(output_folder, f"{Path(ppt_file).stem}.md")
                        save_markdown(markdown_content, output_file)
                        
                st.success("Conversion completed! Check the 'markdown_output' folder.")
        else:
            if folder_path:
                st.error("Invalid folder path. Please enter a valid path.")
    except Exception as e:
        st.error(f"An error occurred: {str(e)}")
        sys.exit(1)

if __name__ == "__main__":
    main()

